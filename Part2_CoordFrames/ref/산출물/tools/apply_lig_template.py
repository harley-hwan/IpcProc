# -*- coding: utf-8 -*-
"""Chapter2 발표자료를 LIG PPT 템플릿에 맞춰 다시 입힌다.

원본 슬라이드의 내용(도형·표·그림·노트)은 그대로 두고
  · 색은 발표자료 원래의 네이비 계열을 그대로 두고, 틀(머리띠·표지·간지·목차)도 같은 남색으로 맞춘다
  · 글꼴 이름을 템플릿과 같은 '맑은 고딕' 으로 통일
  · 표지 / 목차 / 간지 / 본문 페이지의 틀(머리띠·꼬리말·워드마크·쪽번호)을 템플릿 형식으로 교체
  · 목차를 표지 바로 뒤(2장)로 옮긴다
한다.

    python apply_lig_template.py <원본.pptx> <결과.pptx>
"""
import copy
import os
import re
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ───────────────────────────────────────────── 템플릿 디자인 토큰 (네이비판)
# 틀은 LIG 템플릿의 것, 색은 발표자료 원래의 네이비 계열이다
DARK = "13233F"        # 머리띠 · 표지 · 간지 · 목차 왼쪽 패널 바탕 (원본 어두운 장과 같은 남색)
BG = "E8EDF5"          # 본문 슬라이드 바탕 (흰 패널이 얹히는 옅은 청회색)
PANEL = "FFFFFF"       # 흰 패널
TEXT = "1A2233"        # 본문 글자 (원본과 같음)
MUTED = "5A6478"       # 보조 글자 (원본과 같음)
LINE = "B9C0CC"        # 꼬리말 줄 · 가로줄
WHITE = "FFFFFF"
ACCENT = "C2521C"      # 강조 주황 (원본과 같음)
ON_DARK = "C9D6EC"     # 어두운 바탕 위 부제 · 설명 글자
ON_DARK_DIM = "8FA6CC" # 어두운 바탕 위 흐린 글자 (Part 표기 등)
RULE_SOFT = "D0D5DE"   # 목차 행 사이 얇은 줄

KO = "맑은 고딕"
EN = "Arial"
COMPANY = "LIG Defense&Aerospace"
WORDMARK = "LIG"

SW, SH = 13.333, 7.5          # 슬라이드 크기 (in)
M = 0.5                       # 바깥 여백
INNER_W = SW - 2 * M          # 12.333
BAR_Y, BAR_H = 0.45, 0.85     # 본문 머리띠
BAR_BOT = BAR_Y + BAR_H       # 1.30
PANEL_Y, PANEL_H = 1.36, 5.52  # 흰 패널 → 아래끝 6.88
BODY_TOP, BODY_BOT = 1.50, 6.82
FOOT_LINE_Y = 6.95
FOOT_TEXT_Y = 7.00

# ───────────────────────────────────────────────────── 원본 색 → 템플릿 색
# 네이비판은 발표자료 원래 색을 그대로 둔다 (치환 없음).
# 무채색판을 만들 일이 있으면 여기에 {"13233F": "2C2C2C", "2B57A6": "2C2C2C", ...} 식으로 적는다.
PALETTE = {}

# ───────────────────────────────────────── 슬라이드별 역할 (원본 번호 기준)
FIG_DIR = None    # main() 에서 스크립트 위치를 기준으로 채운다
# 이번 작업에서 다시 그린 그림 (슬라이드 번호 → 파일 이름)
REDRAWN = {2: "fig15_scope_vs_chart.png",   # 차콜 배경에 맞춰 색을 무채색으로
           13: "fig06_body.png"}            # 좌우 제목이 겹치던 것을 고침

COVER = 1
HOOK = 2          # 어두운 도입 장
INDEX = 3
DIVIDERS = {9: "1", 22: "2", 29: "3"}
SUMMARY = 34      # 어두운 정리 장
CLOSING = 35

# 본문 슬라이드의 머리띠 번호 (INDEX 의 00~03 과 맞춘다)
def group_of(n):
    if 4 <= n <= 8:
        return "00"
    if 10 <= n <= 21:
        return "01"
    if 23 <= n <= 28:
        return "02"
    if 30 <= n <= 33:
        return "03"
    return None


# ──────────────────────────────────────────────────────────── XML 도우미
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def recolor_xml(el):
    """요소 아래 모든 srgbClr 값을 팔레트에 따라 바꾼다."""
    for c in el.iter(qn("a:srgbClr")):
        v = (c.get("val") or "").upper()
        if v in PALETTE:
            c.set("val", PALETTE[v])


def retypeface(el):
    for t in el.iter():
        tf = t.get("typeface")
        if tf == "Malgun Gothic":
            t.set("typeface", KO)


def set_background(slide, rgb):
    from pptx.oxml import parse_xml
    csld = slide._element.find(qn("p:cSld"))
    old = csld.find(qn("p:bg"))
    if old is not None:
        csld.remove(old)
    xml = (
        f'<p:bg xmlns:p="{P}" xmlns:a="{A}"><p:bgPr>'
        f'<a:solidFill><a:srgbClr val="{rgb}"/></a:solidFill>'
        f"<a:effectLst/></p:bgPr></p:bg>"
    )
    csld.insert(0, parse_xml(xml))


def send_to_back(shape, spTree):
    el = shape._element
    spTree.remove(el)
    spTree.insert(2, el)          # nvGrpSpPr, grpSpPr 다음


def drop(shape):
    shape._element.getparent().remove(shape._element)


def swap_picture(slide, path):
    """슬라이드의 첫 그림을 다른 파일로 갈아끼운다 (가로세로 비율이 같아야 한다)."""
    with open(path, "rb") as f:
        blob = f.read()
    for sh in slide.shapes:
        if sh.shape_type is not None and int(sh.shape_type) == 13:
            part = sh.part.related_part(sh._element.blip_rId)
            part._blob = blob
            return True
    return False


# ──────────────────────────────────────────────────── 도형 만들기 도우미
def txt(slide, text, x, y, w, h, size=12, color=TEXT, bold=False, font=KO,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, spacing=None,
        space_after=0):
    """text 는 문자열 또는 (문자열, 옵션dict) 목록 — 목록이면 문단마다 다른 서식."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = text if isinstance(text, list) else [(text, {})]
    for i, item in enumerate(items):
        s, opt = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opt.get("align", align)
        ls = opt.get("spacing", spacing)
        if ls:
            p.line_spacing = Pt(ls)
        sa = opt.get("space_after", space_after)
        if sa:
            p.space_after = Pt(sa)
        r = p.add_run()
        r.text = s
        f = r.font
        f.name = opt.get("font", font)
        f.size = Pt(opt.get("size", size))
        f.bold = opt.get("bold", bold)
        f.color.rgb = RGBColor.from_string(opt.get("color", color))
    return box


def rect(slide, x, y, w, h, fill, line=None, lw=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    sp.text_frame.word_wrap = False
    return sp


def hline(slide, x1, x2, y, color=LINE, lw=0.75):
    from pptx.enum.shapes import MSO_CONNECTOR
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),
                                    Inches(y), Inches(x2), Inches(y))
    ln.line.color.rgb = RGBColor.from_string(color)
    ln.line.width = Pt(lw)
    return ln


SLIDENUM_XML = (
    '<a:fld xmlns:a="%s" id="{1D0B0B0B-0000-4000-8000-00000000C0DE}" type="slidenum">'
    "<a:rPr lang=\"ko-KR\" altLang=\"en-US\" sz=\"%d\" b=\"0\" dirty=\"0\">"
    '<a:solidFill><a:srgbClr val="%s"/></a:solidFill>'
    '<a:latin typeface="%s"/><a:ea typeface="%s"/></a:rPr>'
    "<a:t>1</a:t></a:fld>"
)


def slide_number(slide, color=TEXT, size=10, x=6.2, y=FOOT_TEXT_Y):
    from pptx.oxml import parse_xml
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.0), Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p._p.append(parse_xml(SLIDENUM_XML % (A, int(size * 100), color, EN, EN)))
    return box


# ─────────────────────────────────────────────────────── 페이지 공통 요소
def content_chrome(slide, title, subtitle):
    """본문 페이지: 어두운 머리띠 + 흰 패널 + 꼬리말."""
    rect(slide, M, BAR_Y, INNER_W, BAR_H, DARK)
    paras = [(title, {"size": 19 if subtitle else 20, "bold": True, "color": WHITE,
                      "space_after": 3})]
    for i, s in enumerate(subtitle):
        paras.append((s, {"size": 11, "color": ON_DARK, "spacing": 14}))
    txt(slide, paras, M + 0.20, BAR_Y, 9.55, BAR_H, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, WORDMARK, 11.15, BAR_Y, 1.50, BAR_H, size=24, bold=True,
        color=WHITE, font=EN, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        wrap=False)
    footer(slide, TEXT)


def dark_chrome(slide, wordmark_size=34):
    """표지·간지 계열: 워드마크 + 가로줄 + 줄 위 회사명."""
    txt(slide, WORDMARK, 0.55, 0.42, 1.6, 0.70, size=wordmark_size, bold=True,
        color=WHITE, font=EN, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    hline(slide, 2.00, 12.83, 0.86, LINE, 0.75)
    box = txt(slide, COMPANY, 10.00, 0.70, 2.40, 0.32, size=11, bold=True,
              color=WHITE, font=EN, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    box.fill.solid()                      # 글자 뒤로 줄이 지나가지 않도록
    box.fill.fore_color.rgb = RGBColor.from_string(DARK)


def footer(slide, color=TEXT, page=True):
    hline(slide, M, SW - M, FOOT_LINE_Y, LINE, 0.75)
    txt(slide, COMPANY, M, FOOT_TEXT_Y, 3.5, 0.30, size=10, bold=True,
        color=color, font=EN, wrap=False)
    if page:
        slide_number(slide, color)


# ────────────────────────────────────────────── 체인(측정▸…▸LLA) 다루기
CHIP_W, ARROW_W, CHIP_H, CHIP_GAP = 1.42, 0.24, 0.44, 0.0


def chain_shapes(slide):
    """표지·간지에 있는 체인 칩/화살표 도형을 찾아 (도형, dx) 목록으로 준다."""
    out = []
    for sh in slide.shapes:
        if sh.width is None or sh.height is None:
            continue
        w, h = sh.width / 914400, sh.height / 914400
        if abs(h - CHIP_H) < 0.02 and (abs(w - CHIP_W) < 0.02 or abs(w - ARROW_W) < 0.02):
            out.append(sh)
    return out


def move_chain(shapes, new_left, new_top):
    if not shapes:
        return
    x0 = min(s.left for s in shapes)
    y0 = min(s.top for s in shapes)
    dx = Inches(new_left) - x0
    dy = Inches(new_top) - y0
    for s in shapes:
        s.left += dx
        s.top += dy


def chain_width(shapes):
    if not shapes:
        return 0
    return (max(s.left + s.width for s in shapes) - min(s.left for s in shapes)) / 914400


# ──────────────────────────────────────────── 본문 내용 세로 위치 다시잡기
def est_text_h(sh):
    """텍스트 상자의 실제 글자 높이 추정 (in)."""
    tf = sh.text_frame
    tot = 0.0
    for p in tf.paragraphs:
        szs = [r.font.size.pt for r in p.runs if r.font.size]
        sz = max(szs) if szs else 12.0
        try:
            ls = p.line_spacing
            line = ls.pt if hasattr(ls, "pt") else (sz * 1.2 * ls if ls else sz * 1.22)
        except Exception:
            line = sz * 1.22
        tot += line
        tot += (p.space_before.pt if p.space_before else 0)
        tot += (p.space_after.pt if p.space_after else 0)
    tot += (tf.margin_top.pt or 0) + (tf.margin_bottom.pt or 0)
    return tot / 72.0


def table_h(tbl, width_in):
    """표는 칸 글이 넘치면 PowerPoint 가 행을 늘리므로, 실제로 그려질 높이를 다시 잰다."""
    tot = 0.0
    ncol = max(1, len(tbl.columns))
    for row in tbl.rows:
        need = 0.0
        for cell, col in zip(row.cells, tbl.columns):
            cw = (col.width / 914400) if col.width else (width_in / ncol)
            mar = ((cell.margin_left.pt if cell.margin_left is not None else 7.2)
                   + (cell.margin_right.pt if cell.margin_right is not None else 7.2)) / 72.0
            inner = max(0.4, cw - mar)
            lines = 0
            sz = 12.0
            for p in cell.text_frame.paragraphs:
                t = p.text
                if not t:
                    lines += 1
                    continue
                szs = [r.font.size.pt for r in p.runs if r.font.size]
                sz = max(szs) if szs else sz
                # 한글은 글자 폭이 글자 크기와 비슷하고, 라틴 문자는 그 절반쯤이다
                w = sum(1.0 if ord(c) > 0x2000 else 0.55 for c in t) * sz / 72.0
                lines += max(1, int(w / inner) + (1 if w % inner else 0))
            h = lines * sz * 1.30 / 72.0 + (
                (cell.margin_top.pt if cell.margin_top is not None else 3.6)
                + (cell.margin_bottom.pt if cell.margin_bottom is not None else 3.6)) / 72.0
            need = max(need, h)
        tot += max((row.height / 914400) if row.height else 0, need)
    return tot


def visual_h(sh):
    dec = sh.height / 914400 if sh.height else 0
    if sh.has_table:
        return max(dec, table_h(sh.table, sh.width / 914400 if sh.width else 12.0))
    if not sh.has_text_frame or not sh.text_frame.text.strip():
        return dec
    anc = sh.text_frame.vertical_anchor
    if anc in (MSO_ANCHOR.MIDDLE, MSO_ANCHOR.BOTTOM):
        return dec
    return min(dec, max(est_text_h(sh), 0.12)) if dec else est_text_h(sh)


def bands_of(shapes):
    """세로로 겹치는 도형끼리 한 '줄'로 묶는다 → [(위, 아래, [도형…]), …]"""
    items = sorted(((s.top / 914400, s.top / 914400 + visual_h(s), s)
                    for s in shapes), key=lambda t: t[0])
    bands = []
    for y0, y1, sh in items:
        if bands and y0 < bands[-1][1] - 0.02:
            bands[-1][1] = max(bands[-1][1], y1)
            bands[-1][2].append(sh)
        else:
            bands.append([y0, y1, [sh]])
    return bands


def refit(shapes, top, bottom, name=""):
    """내용을 [top, bottom] 안으로 옮긴다.

    줄 자체는 절대 줄이지 않고 줄 사이 여백만 균등하게 좁히므로 겹칠 일이 없다.
    여백을 0 으로 줄여도 안 들어가면 경고만 내고 그대로 둔다(그 장은 따로 손본다)."""
    if not shapes:
        return
    bands = bands_of(shapes)
    used = sum(b[1] - b[0] for b in bands)
    gaps = [bands[i + 1][0] - bands[i][1] for i in range(len(bands) - 1)]
    slack = (bottom - top) - used
    if slack < 0:
        # 여백을 다 없애도 안 들어가면 글자 크기까지 같이 줄인다
        s = scale_block(shapes, top, bottom)
        print(f"  · {name}: 통째로 {s:.3f} 배 축소 (줄 높이 합 {used:.2f}in)")
        return
    gs = 1.0
    if gaps and sum(gaps) > 0:
        gs = max(0.0, min(1.0, slack / sum(gaps)))
    cur = top
    for i, (b0, b1, shs) in enumerate(bands):
        dy = cur - b0
        for sh in shs:
            sh.top = Inches(sh.top / 914400 + dy)
        cur += (b1 - b0) + (gaps[i] * gs if i < len(gaps) else 0)


def scale_block(shapes, top, bottom, cx=SW / 2, squeeze_x=False):
    """세로로 통째로 줄여 [top, bottom] 에 맞춘다. 글자 크기도 같이 줄이므로 겹칠 일이 없다.

    가로는 건드리지 않는다 — 좌우로 나란히 둔 것들 사이가 좁아지면 원본에 없던
    겹침이 생기기 때문이다. 다만 그림은 가로세로 비율을 지켜야 하므로
    제자리에서 같은 비율로 줄인다. squeeze_x=True 면 가로도 함께 줄인다."""
    if not shapes:
        return 1.0
    y0 = min(s.top / 914400 for s in shapes)
    y1 = max(s.top / 914400 + visual_h(s) for s in shapes)
    s = min(1.0, (bottom - top) / (y1 - y0)) if y1 > y0 else 1.0
    for sh in shapes:
        x, y, w = sh.left / 914400, sh.top / 914400, (sh.width or 0) / 914400
        sh.top = Inches(top + (y - y0) * s)
        if sh.height:
            sh.height = Inches(sh.height / 914400 * s)
        is_pic = sh.shape_type is not None and int(sh.shape_type) == 13
        if squeeze_x:
            sh.left = Inches(cx + (x - cx) * s)
            if sh.width:
                sh.width = Inches(w * s)
        elif is_pic and sh.width:
            sh.width = Inches(w * s)                 # 비율 유지 — 제자리에서 축소
            sh.left = Inches(x + w * (1 - s) / 2)
        frames = []
        if sh.has_text_frame:
            frames.append(sh.text_frame)
        elif sh.has_table:
            for row in sh.table.rows:
                if row.height:
                    row.height = Inches(row.height / 914400 * s)
                frames.extend(c.text_frame for c in row.cells)
        for tf in frames:
            for p in tf.paragraphs:
                if p.line_spacing is not None and hasattr(p.line_spacing, "pt"):
                    p.line_spacing = Pt(p.line_spacing.pt * s)
                for r in p.runs:
                    if r.font.size:
                        r.font.size = Pt(round(r.font.size.pt * s, 1))
    return s


# ═══════════════════════════════════════════════════════════════ 슬라이드별
def build_cover(slide):
    for sh in list(slide.shapes):          # 체인 칩까지 전부 지운다 — 표지는 템플릿 그대로
        drop(sh)
    set_background(slide, DARK)
    txt(slide, WORDMARK, 0.60, 0.45, 2.2, 1.0, size=54, bold=True, color=WHITE,
        font=EN, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    txt(slide, "Chapter 2", 10.30, 0.55, 2.5, 0.35, size=11, bold=True,
        color=WHITE, font=EN, align=PP_ALIGN.RIGHT, wrap=False)
    txt(slide, "좌표계와 좌표변환", 4.30, 2.70, 8.50, 1.00, size=40, bold=True,
        color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, "안테나가 잰 각도가 지도 위 위경도가 되기까지", 4.30, 3.72, 8.50, 0.42,
        size=16, color=ON_DARK, align=PP_ALIGN.RIGHT)
    txt(slide, COMPANY, 6.80, 4.35, 6.00, 0.40, size=16, bold=True, color=WHITE,
        font=EN, align=PP_ALIGN.RIGHT, wrap=False)
    txt(slide, "레이다 시스템 소프트웨어 스터디  ·  Part 2", 6.80, 4.80, 6.00, 0.35,
        size=12, color=ON_DARK, align=PP_ALIGN.RIGHT)


def build_index(slide):
    for sh in list(slide.shapes):
        drop(sh)
    set_background(slide, BG)
    rect(slide, 0, 0, 4.90, SH, DARK)
    txt(slide, "INDEX", 1.25, 0.55, 2.2, 0.60, size=30, bold=True, color=WHITE,
        font=EN, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    hline(slide, 2.65, 12.00, 0.86, LINE, 0.75)
    txt(slide, WORDMARK, 11.40, 0.40, 1.40, 0.60, size=30, bold=True, color=DARK,
        font=EN, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    txt(slide, "목차", 0.45, 3.05, 4.00, 1.40, size=30, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, COMPANY, 0.45, FOOT_TEXT_Y, 3.5, 0.30, size=10, bold=True,
        color=WHITE, font=EN, wrap=False)
    slide_number(slide, TEXT)

    rows = [
        ("00", "들어가며", "좌표계가 왜 여럿인가 · 변환의 두 동작"),
        ("01", "좌표계 여섯 가지", "안테나 · 동체 · INS\nNED/ENU · ECEF/ECI · LLA"),
        ("02", "실전 설계", "함선 고정형 안테나의 5단계 변환 경로"),
        ("03", "C 로 확인", "구현 · 왕복 시험 · 확인 실험 · UI"),
    ]
    y0, dy = 1.62, 1.22
    for i, (num, name, desc) in enumerate(rows):
        y = y0 + i * dy
        txt(slide, num, 5.55, y, 1.10, 0.62, size=34, bold=True, color=DARK,
            font=EN, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        txt(slide, name, 6.80, y, 2.70, 0.62, size=19, color=TEXT,
            anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        txt(slide, desc, 9.45, y, 3.40, 0.62, size=11.5, color=MUTED,
            anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        if i < len(rows) - 1:
            hline(slide, 5.55, 12.80, y + 0.92, RULE_SOFT, 0.5)
    txt(slide, "사전 지식은 삼각함수와 행렬 곱셈 정도면 충분합니다.  그마저도 필요한 곳에서 다시 설명합니다.",
        5.55, 6.42, 7.30, 0.32, size=11.5, color=MUTED)


def build_divider(slide, part, title, desc, tail):
    chain = chain_shapes(slide)
    for sh in list(slide.shapes):
        if sh not in chain:
            drop(sh)
    set_background(slide, DARK)
    dark_chrome(slide)
    txt(slide, f"Part {part}", 0.95, 2.20, 6.0, 0.42, size=19, bold=True,
        color=ON_DARK_DIM, font=EN, wrap=False)
    txt(slide, title, 0.95, 2.72, 11.5, 0.95, size=40, bold=True, color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, desc, 0.95, 3.82, 11.5, 0.45, size=16, color=ON_DARK)
    if chain:
        move_chain(chain, 0.95, 4.75)
    if tail:
        txt(slide, tail, 0.95, 5.72, 11.5, 0.40, size=14, color=ON_DARK_DIM)
    footer(slide, WHITE)


def build_closing(slide):
    chain = chain_shapes(slide)
    for sh in list(slide.shapes):
        if sh not in chain:
            drop(sh)
    set_background(slide, DARK)
    txt(slide, WORDMARK, 0.60, 0.45, 2.2, 1.0, size=54, bold=True, color=WHITE,
        font=EN, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    txt(slide, "감사합니다", 4.30, 2.70, 8.50, 1.00, size=44, bold=True,
        color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, "질문 환영합니다", 4.30, 3.75, 8.50, 0.42, size=18, color=ON_DARK,
        align=PP_ALIGN.RIGHT)
    txt(slide, COMPANY, 6.80, 4.35, 6.00, 0.40, size=16, bold=True, color=WHITE,
        font=EN, align=PP_ALIGN.RIGHT, wrap=False)
    txt(slide, "코드 · 문서 :  radar-sw-study / Part2_CoordFrames", 6.80, 4.80,
        6.00, 0.35, size=12, color=ON_DARK, align=PP_ALIGN.RIGHT)
    move_chain(chain, 12.73 - chain_width(chain), 5.60)


def build_dark_content(slide, title, top, bottom, title_y=1.22, title_size=28,
                       mode="refit", name=""):
    """어두운 바탕을 유지하는 본문 장(도입 2장, 정리 34장)."""
    keep = []
    for sh in list(slide.shapes):
        y = sh.top / 914400 if sh.top is not None else 99
        if y < 1.05 and sh.has_text_frame:      # 원래 제목은 새로 그린다
            drop(sh)
        else:
            keep.append(sh)
    set_background(slide, DARK)
    dark_chrome(slide)
    txt(slide, title, 0.95, title_y, 11.5, 0.58, size=title_size, bold=True,
        color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if mode == "scale":
        scale_block(keep, top, bottom, squeeze_x=True)
    else:
        refit(keep, top, bottom, name)
    footer(slide, WHITE)


def build_content(slide, num):
    """일반 본문 페이지."""
    title, subtitle, body = None, [], []
    for sh in slide.shapes:
        y = sh.top / 914400 if sh.top is not None else 99
        x = sh.left / 914400 if sh.left is not None else 99
        if sh.has_text_frame and sh.text_frame.text.strip():
            szs = [r.font.size.pt for p in sh.text_frame.paragraphs
                   for r in p.runs if r.font.size]
            mx = max(szs) if szs else 0
            if title is None and y < 0.75 and mx >= 24:
                title = sh
                continue
            if title is not None and not subtitle and 0.80 <= y <= 1.30 and mx <= 16 and x < 1.0:
                subtitle = [p.text for p in sh.text_frame.paragraphs if p.text.strip()]
                drop(sh)
                continue
        body.append(sh)

    title_text = title.text_frame.text.replace("\n", " ").strip() if title else ""
    title_text = re.sub(r"\s{2,}", "  ", title_text)
    if title is not None:
        drop(title)

    set_background(slide, BG)
    refit(body, BODY_TOP, BODY_BOT, f"slide {num}·{title_text[:18]}")

    spTree = slide.shapes._spTree
    panel = rect(slide, M, PANEL_Y, INNER_W, PANEL_H, PANEL)
    send_to_back(panel, spTree)

    content_chrome(slide, f"{num}. {title_text}" if num else title_text, subtitle)


# ═══════════════════════════════════════════════════════════════════ 본체
def main(src, dst):
    global FIG_DIR
    # 그림 폴더는 이 스크립트(산출물/tools/) 위치 기준 — 입력 pptx 가 어디 있든 같다
    FIG_DIR = os.path.normpath(os.path.join(
        os.path.dirname(os.path.abspath(__file__)), "..", "..", "참고 자료", "figures"))
    prs = Presentation(src)

    # 0) 다시 그린 그림을 갈아끼운다
    for n, fn in REDRAWN.items():
        path = os.path.join(FIG_DIR, fn)
        if os.path.exists(path) and swap_picture(prs.slides[n - 1], path):
            print(f"  그림 교체: {n}장 ← {fn}")

    # 1) 색 · 글꼴을 먼저 통째로 치환한다 (표·표 테두리까지 한 번에)
    for slide in prs.slides:
        recolor_xml(slide._element)
        retypeface(slide._element)
        if slide.has_notes_slide:
            retypeface(slide.notes_slide._element)
    for master in prs.slide_masters:
        recolor_xml(master._element)
        retypeface(master._element)
        for layout in master.slide_layouts:
            recolor_xml(layout._element)
            retypeface(layout._element)

    # 2) 페이지 틀을 템플릿 형식으로 다시 만든다
    for i, slide in enumerate(prs.slides, 1):
        if i == COVER:
            build_cover(slide)
        elif i == HOOK:
            build_dark_content(
                slide, "레이다가 잰 값과  전시기에 보일 값은 다르다",
                top=1.92, bottom=6.85, title_y=1.18, title_size=27, mode="scale")
        elif i == INDEX:
            build_index(slide)
        elif i in DIVIDERS:
            part = DIVIDERS[i]
            spec = {
                "1": ("좌표계 여섯 가지",
                      "각각 원점이 어디이고, 축이 어디를 향하고, 무엇에 붙어 있는가", ""),
                "2": ("실전 설계",
                      "함선 고정형 안테나를 놓고 변환 경로를 설계합니다",
                      "신호처리  →  데이터 처리 (좌표변환 + 추적)  →  통제제어 전시기"),
                "3": ("C 로 짜서 확인하기",
                      "좌표변환은 \"적당히 그럴듯한 값\"이 나오기 때문에 눈으로는 검증되지 않습니다",
                      "Visual Studio 2022  ·  C17  ·  경고 0건  ·  왕복 오차 1e−9 m 미만"),
            }[part]
            build_divider(slide, part, *spec)
        elif i == SUMMARY:
            build_dark_content(slide, "정리", top=2.05, bottom=6.82,
                               title_y=1.22, title_size=30, name="정리")
        elif i == CLOSING:
            build_closing(slide)
        else:
            build_content(slide, group_of(i))

    # 3) 목차를 표지 바로 뒤로 옮긴다 (원본 순서는 표지 → 도입 → 목차)
    ids = prs.slides._sldIdLst
    el = list(ids)[INDEX - 1]
    ids.remove(el)
    ids.insert(1, el)

    prs.save(dst)
    print("written:", dst)


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
