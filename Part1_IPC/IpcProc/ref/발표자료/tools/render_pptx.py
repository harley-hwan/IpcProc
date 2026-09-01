# -*- coding: utf-8 -*-
"""pptx -> PDF 렌더러. 이 자료가 쓰는 기능(단색 채움 · 기본 도형 · 텍스트 · 표)만 처리한다."""
import pymupdf as fitz
from pptx import Presentation
from pptx.util import Emu

EMU = 12700.0                      # 1 pt
FONTS = {
    ("맑은 고딕", False): r"C:\Windows\Fonts\malgun.ttf",
    ("맑은 고딕", True):  r"C:\Windows\Fonts\malgunbd.ttf",
    ("Consolas", False):  r"C:\Windows\Fonts\consola.ttf",
    ("Consolas", True):   r"C:\Windows\Fonts\consolab.ttf",
}
ALIAS = {"mg": ("맑은 고딕", False), "mgb": ("맑은 고딕", True),
         "cn": ("Consolas", False), "cnb": ("Consolas", True)}
_cache = {}

def face(name, bold):
    n = "Consolas" if name and "consol" in name.lower() else "맑은 고딕"
    key = (n, bool(bold))
    tag = {("맑은 고딕", False): "mg", ("맑은 고딕", True): "mgb",
           ("Consolas", False): "cn", ("Consolas", True): "cnb"}[key]
    if tag not in _cache:
        _cache[tag] = fitz.Font(fontfile=FONTS[key])
    return tag, _cache[tag]

def pt(v):  return (v or 0) / EMU
def rgb(c):
    try:
        if c and c.type is not None and c.rgb is not None:
            v = c.rgb; return (v[0]/255.0, v[1]/255.0, v[2]/255.0)
    except Exception: pass
    return None

def is_cjk(ch): return ord(ch) > 0x2000

def segment(text, name, bold):
    """폰트에 글리프가 없는 글자는 맑은 고딕으로 떨어뜨린다 (Consolas 에 한글 없음)."""
    pt_, pf = face(name, bold)
    ft_, ff = face("맑은 고딕", bold)
    out, buf, cur = [], "", None
    for ch in text:
        use = (pt_, pf) if (ch in " 	" or pf.has_glyph(ord(ch))) else (ft_, ff)
        if cur is not None and use[0] != cur[0]:
            out.append((buf, cur[0], cur[1])); buf = ""
        buf += ch; cur = use
    if buf: out.append((buf, cur[0], cur[1]))
    return out

def wrap(spans, width):
    """spans=[(text,tag,font,size,bold,color)] -> [[span조각,...], ...]"""
    lines, cur, w = [], [], 0.0
    for text, tag, fnt, size, bold, col in spans:
        tok, buf = [], ""
        for ch in text:                      # 공백 단위 + CJK 는 글자 단위로 분해
            if ch == " ": tok.append(buf + " "); buf = ""
            elif is_cjk(ch):
                if buf: tok.append(buf); buf = ""
                tok.append(ch)
            else: buf += ch
        if buf: tok.append(buf)
        for t in tok:
            tw = fnt.text_length(t, size)
            if cur and w + tw > width and t.strip():
                lines.append(cur); cur, w = [], 0.0
            cur.append((t, tag, fnt, size, bold, col)); w += tw
    lines.append(cur)
    return lines

# ---------------- 도형 ----------------
def draw_shape(pg, sh):
    x0, y0 = pt(sh.left), pt(sh.top)
    w, h = pt(sh.width), pt(sh.height)
    if w <= 0 or h <= 0: return
    r = fitz.Rect(x0, y0, x0 + w, y0 + h)
    fill = line = None; lw = 0.0
    try:
        if sh.fill.type == 1: fill = rgb(sh.fill.fore_color)
    except Exception: pass
    try:
        line = rgb(sh.line.color)
        lw = pt(sh.line.width) if sh.line.width else (0.75 if line else 0)
    except Exception: pass
    if fill is None and line is None: return
    try: prst = sh._element.spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom').get('prst')
    except Exception: prst = 'rect'
    s = pg.new_shape()
    if prst == 'roundRect':
        try: s.draw_rect(r, radius=0.16)
        except TypeError: s.draw_rect(r)
    elif prst == 'ellipse': s.draw_oval(r)
    elif prst == 'line':    s.draw_line(fitz.Point(r.x0, r.y0), fitz.Point(r.x1, r.y1))
    elif prst in ('downArrow','upArrow','leftArrow','rightArrow','upDownArrow'):
        a, b = 0.4, 0.45           # 화살대 폭 비율, 머리 길이 비율
        cx, cy = (r.x0+r.x1)/2, (r.y0+r.y1)/2
        if prst in ('downArrow','upArrow','upDownArrow'):
            hw, hh = w*a/2, h*b
            top, bot = (r.y0, r.y1) if prst != 'upArrow' else (r.y1, r.y0)
            pts = [(cx-hw, top), (cx-hw, bot-(hh if prst!='upArrow' else -hh)),
                   (r.x0, bot-(hh if prst!='upArrow' else -hh)), (cx, bot),
                   (r.x1, bot-(hh if prst!='upArrow' else -hh)),
                   (cx+hw, bot-(hh if prst!='upArrow' else -hh)), (cx+hw, top)]
        else:
            hh, hw = h*a/2, w*b
            lft, rgt = (r.x0, r.x1) if prst == 'rightArrow' else (r.x1, r.x0)
            sgn = 1 if prst == 'rightArrow' else -1
            pts = [(lft, cy-hh), (rgt-sgn*hw, cy-hh), (rgt-sgn*hw, r.y0), (rgt, cy),
                   (rgt-sgn*hw, r.y1), (rgt-sgn*hw, cy+hh), (lft, cy+hh)]
        s.draw_polyline([fitz.Point(*p) for p in pts] + [fitz.Point(*pts[0])])
    else: s.draw_rect(r)
    s.finish(fill=fill, color=line, width=lw or 0, closePath=True)
    s.commit()

AN = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

def para_meta(p):
    """(marL, indent, 글머리기호) — pptx 가 노출 안 하는 값이라 XML 에서 직접 읽음"""
    pr = p._p.find(AN + "pPr")
    if pr is None: return 0.0, 0.0, None
    marL = float(pr.get("marL") or 0) / EMU
    ind  = float(pr.get("indent") or 0) / EMU
    ch = None
    if pr.find(AN + "buNone") is None:
        b = pr.find(AN + "buChar")
        if b is not None: ch = b.get("char")
    return marL, ind, ch

# ---------------- 텍스트 ----------------
def draw_text(pg, sh, fonts):
    tf = sh.text_frame
    if not tf.text.strip(): return
    ml = pt(tf.margin_left) if tf.margin_left is not None else 7.2
    mr = pt(tf.margin_right) if tf.margin_right is not None else 7.2
    mt = pt(tf.margin_top) if tf.margin_top is not None else 3.6
    mb = pt(tf.margin_bottom) if tf.margin_bottom is not None else 3.6
    x0, y0 = pt(sh.left) + ml, pt(sh.top) + mt
    aw = pt(sh.width) - ml - mr
    ah = pt(sh.height) - mt - mb
    if aw <= 1: return
    blocks = []
    for p in tf.paragraphs:
        spans = []
        for run in p.runs:
            size = run.font.size.pt if run.font.size else (p.font.size.pt if p.font.size else 18.0)
            bold = run.font.bold if run.font.bold is not None else bool(p.font.bold)
            col = rgb(run.font.color) or rgb(p.font.color) or (0, 0, 0)
            for txt, tag, fnt in segment(run.text, run.font.name or p.font.name, bold):
                if tag not in fonts:
                    pg.insert_font(fontname=tag, fontfile=FONTS[ALIAS[tag]]); fonts.add(tag)
                spans.append((txt, tag, fnt, size, bold, col))
        base = max([s[3] for s in spans], default=12.0)
        ls = p.line_spacing
        if ls is None:            lh = base * 1.2
        elif isinstance(ls, float): lh = ls * base * 1.2
        else:                     lh = pt(ls)          # Length = 줄 높이 그 자체
        marL, ind, bu = para_meta(p)
        lines = wrap(spans, aw - marL) if spans else [[]]
        al = {2: "ctr", 3: "r"}.get(int(p.alignment) if p.alignment is not None else 1, "l")
        col0 = spans[0][5] if spans else (0, 0, 0)
        blocks.append((lines, lh, al, base, marL, ind, bu, col0))
    total = sum(len(b[0]) * b[1] for b in blocks)
    anc = int(tf.vertical_anchor) if tf.vertical_anchor is not None else 1
    y = y0 + (max(0, ah - total) / 2 if anc == 3 else (max(0, ah - total) if anc == 4 else 0))
    for lines, lh, al, base, marL, ind, bu, col0 in blocks:
        tw_ = aw - marL
        for li, ln in enumerate(lines):
            wsum = sum(f.text_length(t, s) for t, _, f, s, _, _ in ln)
            xl = x0 + marL
            x = xl + (tw_ - wsum) / 2 if al == "ctr" else (xl + tw_ - wsum if al == "r" else xl)
            y += lh
            by = y - lh * 0.28
            if bu and li == 0 and any(t.strip() for t, *_ in ln):
                btag, bfnt = face("맑은 고딕", False)
                if btag not in fonts:
                    pg.insert_font(fontname=btag, fontfile=FONTS[ALIAS[btag]]); fonts.add(btag)
                pg.insert_text(fitz.Point(xl + ind, by), bu, fontname=btag,
                               fontsize=base * 0.9, color=col0)
            # 같은 스타일끼리 이어 붙여 한 번에 출력 (공백이 PDF 에 남아야 복사가 정상)
            run = []
            for t, tag, f, s_, b, c in ln:
                if run and (run[-1][1], run[-1][3], run[-1][5]) == (tag, s_, c):
                    run[-1] = (run[-1][0] + t, tag, f, s_, b, c)
                else:
                    run.append((t, tag, f, s_, b, c))
            for t, tag, f, s_, b, c in run:
                pg.insert_text(fitz.Point(x, by), t, fontname=tag, fontsize=s_, color=c)
                x += f.text_length(t, s_)

# ---------------- 표 ----------------
def draw_table(pg, gf, fonts):
    tb = gf.table
    ys, y = [], pt(gf.top)
    for r in tb.rows: ys.append(y); y += pt(r.height)
    ys.append(y)
    xs, x = [], pt(gf.left)
    for c in tb.columns: xs.append(x); x += pt(c.width)
    xs.append(x)
    for ri, row in enumerate(tb.rows):
        for ci, cell in enumerate(row.cells):
            r = fitz.Rect(xs[ci], ys[ri], xs[ci+1], ys[ri+1])
            f = None
            try:
                if cell.fill.type == 1: f = rgb(cell.fill.fore_color)
            except Exception: pass
            s = pg.new_shape(); s.draw_rect(r)
            s.finish(fill=f, color=(0.75, 0.75, 0.75), width=0.5); s.commit()
            if cell.text.strip():
                class P:  # draw_text 가 기대하는 최소 인터페이스
                    left, top, width, height = int(r.x0*EMU), int(r.y0*EMU), int(r.width*EMU), int(r.height*EMU)
                    text_frame = cell.text_frame
                draw_text(pg, P, fonts)

# ---------------- 메인 ----------------
def walk(pg, shapes, fonts):
    for sh in shapes:
        try:
            if sh.shape_type == 6:            # GROUP
                walk(pg, sh.shapes, fonts); continue
            if getattr(sh, "has_table", False) and sh.has_table:
                draw_table(pg, sh, fonts); continue
            if sh.shape_type != 17:           # TEXT_BOX 는 도형 그리기 생략 대상 아님(채움 있을 수 있음)
                draw_shape(pg, sh)
            else:
                draw_shape(pg, sh)
            if sh.has_text_frame: draw_text(pg, sh, fonts)
        except Exception as e:
            print("   ! %s: %s" % (getattr(sh, "name", "?"), e))

def render(src, out):
    prs = Presentation(src)
    W, H = pt(prs.slide_width), pt(prs.slide_height)
    doc = fitz.open()
    for i, sl in enumerate(prs.slides, 1):
        pg = doc.new_page(width=W, height=H); fonts = set()
        # 슬라이드마다 p:bg 로 배경색이 지정돼 있다 (표지는 어두운 색)
        col = (1, 1, 1)
        el = sl._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
        if el is not None:
            c = el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if c is not None:
                v = c.get('val'); col = tuple(int(v[i:i+2], 16) / 255.0 for i in (0, 2, 4))
        bg = pg.new_shape(); bg.draw_rect(fitz.Rect(0, 0, W, H))
        bg.finish(fill=col, color=None); bg.commit()
        walk(pg, sl.slide_layout.shapes, fonts)   # 레이아웃 배경 먼저
        walk(pg, sl.shapes, fonts)
    doc.set_metadata({"title": "IPC (Inter Process Communication)", "producer": "IpcProc render_pptx.py"})
    doc.subset_fonts(verbose=False)      # 쓴 글자만 남김 (안 하면 한글 폰트 통째로 들어가 수십 MB)
    doc.save(out, deflate=True, garbage=4); doc.close()
    print("완료: %s (%d쪽, %.0fx%.0f)" % (out, len(prs.slides.__iter__().__length_hint__() * [0]) if False else len(prs.slides._sldIdLst), W, H))

if __name__ == "__main__":
    import sys
    render(sys.argv[1], sys.argv[2])
